#!/usr/bin/env python3
"""
sync_docs.py — keep the ticketing-project artifacts in lockstep with the brief.

The brief (docs/ticketing-system-brief.md) is the single source of truth. This
helper does the parts that are tedious or error-prone to do by hand:

  inspect        Print every scope-relevant value across all artifacts, side by side.
  verify         Diff each artifact against the brief; non-zero exit on any drift.
  parse-brief    Emit the canonical FR list + MoSCoW counts (JSON) from the brief.
  xlsx-set-scope Write MUST/SHOULD/COULD (+ optional FUTURE) counts to the workbook.
  xlsx-upsert-fr Insert or update one FR row in the Requirements sheet.
  pptx-replace   Find/replace text on a given slide (for the Scope slide, #6).

Markdown and HTML are plain text — edit those with your normal text tools; this
helper only owns the binary docs (.xlsx/.pptx) plus the read-only verify/inspect
checks that span every artifact. Project memory lives outside the repo and is
handled by the calling command, not here.

Run from anywhere: paths resolve relative to this file's repo root.
"""
from __future__ import annotations

import argparse
import json
import re
import sys
from pathlib import Path

ROOT = Path(__file__).resolve().parent.parent
DOCS = ROOT / "docs"
BRIEF = DOCS / "ticketing-system-brief.md"
HTML = DOCS / "ticketing-system.html"
XLSX = DOCS / "Ticketing-Requirements.xlsx"
PPTX = DOCS / "Ticketing-Kickoff.pptx"
CLAUDEMD = ROOT / "CLAUDE.md"

PRIORITIES = ("MUST", "SHOULD", "COULD")
SCOPE_SHEET = "Scope Summary"
REQ_SHEET = "Requirements"


# ----------------------------------------------------------------------------
# Source of truth: parse the brief
# ----------------------------------------------------------------------------
def parse_brief() -> dict:
    """Return {'frs': {id: {'priority','text'}}, 'counts': {MUST,SHOULD,COULD}}."""
    if not BRIEF.exists():
        die(f"brief not found: {BRIEF}")
    text = BRIEF.read_text(encoding="utf-8")
    frs: dict[str, dict] = {}
    # Lines like:  - FR-7: A client can raise a ticket ... *(MUST)*
    pat = re.compile(r"^\s*-\s*\*?\*?FR-(\d+)\*?\*?:\s*(.+?)\s*\*\((MUST|SHOULD|COULD)\)\*",
                     re.MULTILINE)
    for m in pat.finditer(text):
        fid = f"FR-{m.group(1)}"
        body = re.sub(r"\*\*|\*", "", m.group(2)).strip()  # strip md emphasis
        frs[fid] = {"priority": m.group(3), "text": body}
    counts = {p: sum(1 for v in frs.values() if v["priority"] == p) for p in PRIORITIES}
    return {"frs": frs, "counts": counts}


# ----------------------------------------------------------------------------
# Extractors: pull the same numbers back out of each artifact
# ----------------------------------------------------------------------------
def html_values() -> dict:
    """Counts + FR set embedded in the HTML one-pager."""
    if not HTML.exists():
        return {}
    t = HTML.read_text(encoding="utf-8")
    out: dict = {}
    # Legend rows: <span class="cnt">17</span> in order MUST/SHOULD/COULD
    legend = re.findall(r'class="cnt">(\d+)<', t)
    if len(legend) >= 3:
        out["legend_counts"] = dict(zip(PRIORITIES, (int(x) for x in legend[:3])))
    # Filter buttons: data-f="MUST">Must <span class="b">17</span>
    btn = {}
    for p in PRIORITIES:
        b = re.search(rf'data-f="{p}">[^<]*<span class="b">(\d+)<', t)
        if b:
            btn[p] = int(b.group(1))
    if btn:
        out["button_counts"] = btn
    # Donut centre total + "All <n>" + "N concrete requirements"
    donut = re.search(r'class="big">(\d+)<', t)
    if donut:
        out["donut_total"] = int(donut.group(1))
    allbtn = re.search(r'data-f="all">[^<]*<span class="b">(\d+)<', t)
    if allbtn:
        out["all_count"] = int(allbtn.group(1))
    # REQS JS array -> {id: priority}
    reqs = {}
    for m in re.finditer(r'\["(FR-\d+)","(?:[^"\\]|\\.)*?","[^"]*?","(MUST|SHOULD|COULD)"\]', t):
        reqs[m.group(1)] = m.group(2)
    out["reqs"] = reqs
    return out


def xlsx_values() -> dict:
    import openpyxl
    if not XLSX.exists():
        return {}
    wb = openpyxl.load_workbook(XLSX, data_only=True)
    out: dict = {}
    if SCOPE_SHEET in wb.sheetnames:
        ws = wb[SCOPE_SHEET]
        counts = {}
        for row in ws.iter_rows(values_only=True):
            if row and row[0] in PRIORITIES:
                counts[row[0]] = _to_int(row[-1])
        out["scope_counts"] = counts
    if REQ_SHEET in wb.sheetnames:
        ws = wb[REQ_SHEET]
        reqs = {}
        header_row, cols = _find_req_header(ws)
        if header_row:
            for row in ws.iter_rows(min_row=header_row + 1, values_only=True):
                fid = row[cols["ID"]] if cols["ID"] < len(row) else None
                if fid and str(fid).startswith("FR-"):
                    pr = row[cols["Priority"]] if cols["Priority"] < len(row) else None
                    reqs[str(fid)] = pr_norm(pr)
        out["reqs"] = reqs
    return out


def claudemd_values() -> dict:
    if not CLAUDEMD.exists():
        return {}
    t = CLAUDEMD.read_text(encoding="utf-8")
    out = {}
    counts = {}
    for p in PRIORITIES:
        m = re.search(rf"{p}\s*\((\d+)", t)  # e.g. "COULD (4, do not commit)"
        if m:
            counts[p] = int(m.group(1))
    if counts:
        out["scope_counts"] = counts
    tbl = re.search(r"\((\d+)\s*tables\)", t)
    if tbl:
        out["tables"] = int(tbl.group(1))
    return out


def pptx_slide6_text() -> str:
    from pptx import Presentation
    if not PPTX.exists():
        return ""
    prs = Presentation(str(PPTX))
    slides = list(prs.slides)
    if len(slides) < 6:
        return ""
    chunks = []
    for sh in slides[5].shapes:
        if sh.has_text_frame and sh.text.strip():
            chunks.append(sh.text.strip())
    return "\n".join(chunks)


# ----------------------------------------------------------------------------
# Commands
# ----------------------------------------------------------------------------
def cmd_parse_brief(_args) -> int:
    print(json.dumps(parse_brief(), indent=2, ensure_ascii=False))
    return 0


def cmd_inspect(_args) -> int:
    b = parse_brief()
    print("== BRIEF (source of truth) ==")
    print(f"   FRs: {len(b['frs'])}   counts: {fmt_counts(b['counts'])}")
    print("\n== HTML (ticketing-system.html) ==")
    h = html_values()
    print(f"   legend:  {fmt_counts(h.get('legend_counts'))}")
    print(f"   buttons: {fmt_counts(h.get('button_counts'))}")
    print(f"   donut total: {h.get('donut_total')}   all-btn: {h.get('all_count')}   REQS rows: {len(h.get('reqs', {}))}")
    print("\n== XLSX (Ticketing-Requirements.xlsx) ==")
    x = xlsx_values()
    print(f"   Scope Summary: {fmt_counts(x.get('scope_counts'))}")
    print(f"   Requirements rows: {len(x.get('reqs', {}))}")
    print("\n== CLAUDE.md ==")
    c = claudemd_values()
    print(f"   scope: {fmt_counts(c.get('scope_counts'))}   data model: {c.get('tables')} tables")
    print("\n== PPTX slide 6 (Scope) ==")
    print(indent(pptx_slide6_text() or "(empty / not found)"))
    return 0


def cmd_verify(_args) -> int:
    b = parse_brief()
    truth = b["counts"]
    truth_frs = {k: v["priority"] for k, v in b["frs"].items()}
    problems: list[str] = []

    def check_counts(label, got):
        if not got:
            problems.append(f"{label}: no counts found (could not extract)")
            return
        for p in PRIORITIES:
            if got.get(p) != truth[p]:
                problems.append(f"{label}: {p}={got.get(p)} but brief={truth[p]}")

    def check_frs(label, got):
        if not got:
            problems.append(f"{label}: no FR set found")
            return
        miss = sorted(set(truth_frs) - set(got), key=fr_key)
        extra = sorted(set(got) - set(truth_frs), key=fr_key)
        if miss:
            problems.append(f"{label}: missing FRs {miss}")
        if extra:
            problems.append(f"{label}: unknown FRs {extra}")
        for fid in sorted(set(truth_frs) & set(got), key=fr_key):
            if got[fid] != truth_frs[fid]:
                problems.append(f"{label}: {fid} priority {got[fid]} but brief={truth_frs[fid]}")

    h = html_values()
    check_counts("HTML legend", h.get("legend_counts"))
    check_counts("HTML buttons", h.get("button_counts"))
    check_frs("HTML REQS", h.get("reqs"))
    total = sum(truth.values())
    if h.get("donut_total") not in (None, total):
        problems.append(f"HTML donut total={h.get('donut_total')} but brief total={total}")
    if h.get("all_count") not in (None, total):
        problems.append(f"HTML all-button={h.get('all_count')} but brief total={total}")

    x = xlsx_values()
    check_counts("XLSX Scope Summary", x.get("scope_counts"))
    check_frs("XLSX Requirements", x.get("reqs"))

    c = claudemd_values()
    check_counts("CLAUDE.md", c.get("scope_counts"))

    print(f"Brief: {len(b['frs'])} FRs  {fmt_counts(truth)}  total={total}\n")
    if problems:
        print(f"DRIFT — {len(problems)} mismatch(es):")
        for p in problems:
            print(f"  ✗ {p}")
        print("\nFAIL")
        return 1
    print("All artifacts consistent with the brief.  ✓\nOK")
    return 0


def cmd_xlsx_set_scope(args) -> int:
    import openpyxl
    wb = openpyxl.load_workbook(XLSX)
    if SCOPE_SHEET not in wb.sheetnames:
        die(f"sheet '{SCOPE_SHEET}' not found")
    ws = wb[SCOPE_SHEET]
    want = {"MUST": args.must, "SHOULD": args.should, "COULD": args.could}
    last_label_row = None
    for row in ws.iter_rows():
        label = row[0].value
        if label in want and want[label] is not None:
            row[-1].value = want[label]
            last_label_row = row[0].row
        if label == "TOTAL":
            total = sum(v for v in want.values() if v is not None)
            row[-1].value = total
    if args.future is not None and last_label_row is not None:
        # Add/update a FUTURE row just above TOTAL if not present.
        existing = None
        for row in ws.iter_rows():
            if row[0].value in ("FUTURE", "FUTURE / WON'T", "WON'T"):
                existing = row
                break
        if existing:
            existing[-1].value = args.future
        else:
            print("note: no FUTURE row present; add it manually to preserve sheet styling")
    wb.save(XLSX)
    print(f"xlsx scope updated: {fmt_counts(want)}")
    return 0


def cmd_xlsx_upsert_fr(args) -> int:
    import openpyxl
    wb = openpyxl.load_workbook(XLSX)
    ws = wb[REQ_SHEET]
    header_row, cols = _find_req_header(ws)
    if not header_row:
        die("could not locate Requirements header row")
    target = None
    for r in range(header_row + 1, ws.max_row + 1):
        if str(ws.cell(r, cols["ID"] + 1).value) == args.id:
            target = r
            break
    if target is None:
        target = ws.max_row + 1
        ws.cell(target, cols["ID"] + 1).value = args.id
    if args.category is not None:
        ws.cell(target, cols["Category"] + 1).value = args.category
    if args.text is not None:
        ws.cell(target, cols["Requirement"] + 1).value = args.text
    if args.priority is not None:
        ws.cell(target, cols["Priority"] + 1).value = args.priority
    wb.save(XLSX)
    print(f"xlsx {args.id} upserted at row {target}")
    return 0


def cmd_pptx_replace(args) -> int:
    from pptx import Presentation
    prs = Presentation(str(PPTX))
    slides = list(prs.slides)
    if args.slide < 1 or args.slide > len(slides):
        die(f"slide {args.slide} out of range (1..{len(slides)})")
    hits = 0
    for sh in slides[args.slide - 1].shapes:
        if not sh.has_text_frame:
            continue
        for para in sh.text_frame.paragraphs:
            runs = para.runs
            joined = "".join(r.text for r in runs)
            if args.find in joined:
                new = joined.replace(args.find, args.replace)
                if runs:
                    runs[0].text = new
                    for r in runs[1:]:
                        r.text = ""
                    hits += 1
    if hits == 0:
        die(f"text {args.find!r} not found on slide {args.slide}")
    prs.save(str(PPTX))
    print(f"pptx slide {args.slide}: {hits} replacement(s)")
    return 0


# ----------------------------------------------------------------------------
# Helpers
# ----------------------------------------------------------------------------
def _find_req_header(ws):
    """Find the 'ID | Category | Requirement | ... | Priority' header row (0-based col idx)."""
    for row in ws.iter_rows(min_row=1, max_row=8):
        vals = [str(c.value).strip() if c.value is not None else "" for c in row]
        if "ID" in vals and "Priority" in vals and "Requirement" in vals:
            cols = {name: vals.index(name) for name in ("ID", "Category", "Requirement", "Priority")}
            return row[0].row, cols
    return None, None


def _to_int(v):
    try:
        return int(v)
    except (TypeError, ValueError):
        return None


def pr_norm(v):
    s = (str(v) or "").strip().upper()
    return s if s in PRIORITIES else s


def fr_key(fid: str) -> int:
    m = re.search(r"\d+", fid)
    return int(m.group()) if m else 0


def fmt_counts(d) -> str:
    if not d:
        return "(none)"
    return "  ".join(f"{p}={d.get(p)}" for p in PRIORITIES)


def indent(s: str, n: int = 3) -> str:
    pad = " " * n
    return "\n".join(pad + ln for ln in s.splitlines())


def die(msg: str):
    print(f"ERROR: {msg}", file=sys.stderr)
    sys.exit(2)


def main() -> int:
    # Windows consoles default to cp1252; force UTF-8 so ✓/✗ and the docs' em
    # dashes / smart quotes print instead of crashing.
    for stream in (sys.stdout, sys.stderr):
        try:
            stream.reconfigure(encoding="utf-8")
        except (AttributeError, ValueError):
            pass
    ap = argparse.ArgumentParser(description="Sync ticketing-project docs against the brief.")
    sub = ap.add_subparsers(dest="cmd", required=True)
    sub.add_parser("inspect", help="show every scope value across all artifacts")
    sub.add_parser("verify", help="diff all artifacts against the brief (exit 1 on drift)")
    sub.add_parser("parse-brief", help="emit canonical FR list + counts as JSON")

    s = sub.add_parser("xlsx-set-scope", help="write scope counts to the workbook")
    s.add_argument("--must", type=int)
    s.add_argument("--should", type=int)
    s.add_argument("--could", type=int)
    s.add_argument("--future", type=int, default=None)

    s = sub.add_parser("xlsx-upsert-fr", help="insert/update one FR row")
    s.add_argument("--id", required=True)
    s.add_argument("--priority")
    s.add_argument("--category")
    s.add_argument("--text")

    s = sub.add_parser("pptx-replace", help="find/replace text on a slide")
    s.add_argument("--slide", type=int, required=True)
    s.add_argument("--find", required=True)
    s.add_argument("--replace", required=True)

    args = ap.parse_args()
    return {
        "inspect": cmd_inspect,
        "verify": cmd_verify,
        "parse-brief": cmd_parse_brief,
        "xlsx-set-scope": cmd_xlsx_set_scope,
        "xlsx-upsert-fr": cmd_xlsx_upsert_fr,
        "pptx-replace": cmd_pptx_replace,
    }[args.cmd](args)


if __name__ == "__main__":
    sys.exit(main())
